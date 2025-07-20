import os
import tempfile
from typing import List

import fitz  # PyMuPDF
import requests
from pptx import Presentation

max_page_count = 100


def process_file(
    uploaded_url: str, page_numbers: List[int]
) -> List[str]:
    try:
        response = requests.get(uploaded_url)
        file_content = response.content

        if uploaded_url.endswith(".pdf"):
            pdf_documents = fitz.open(stream=file_content, filetype="pdf")

            if max_page_count < len(pdf_documents):
                raise ValueError(f"페이지 수가 {max_page_count}페이지를 초과합니다.")

            one_based_pages = [""]
            for pdf_document in pdf_documents:
                one_based_pages.append(pdf_document.get_text())
            pdf_documents.close()

            if not page_numbers:
                return one_based_pages

            select_pages = [
                one_based_pages[i] for i in page_numbers if 0 < i < len(one_based_pages)
            ]
            return [""] + select_pages

        elif uploaded_url.endswith(".pptx"):
            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name

            presentation = Presentation(temp_file_path)
            if max_page_count < len(presentation.slides):
                raise ValueError(f"페이지 수가 {max_page_count}페이지를 초과합니다.")

            one_based_pages = [""]
            for slide in presentation.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n\n"
                one_based_pages.append(slide_text)

            os.unlink(temp_file_path)

            if not page_numbers:
                return one_based_pages

            select_pages = [
                one_based_pages[i] for i in page_numbers if 0 < i < len(one_based_pages)
            ]
            return [""] + select_pages
        else:
            raise ValueError("지원하지 않는 파일 형식입니다.")
    except Exception as e:
        raise e
